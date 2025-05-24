import pytest
import tempfile
import os
from document_processor import DocumentProcessor, DocumentChunk
from pptx import Presentation
from docx import Document
import PyPDF2

class TestDocumentProcessor:
    """Test document processing functionality"""
    
    @pytest.fixture
    def processor(self):
        return DocumentProcessor(chunk_size=100, chunk_overlap=20)
    
    def test_chunk_creation(self, processor):
        """Test text chunking"""
        text = " ".join([f"word{i}" for i in range(250)])  # 250 words
        chunks = processor._create_chunks(text, page_number=1, chunk_type="test")
        
        # Should create 3 chunks with overlap
        assert len(chunks) >= 3
        
        # Check chunk properties
        for i, chunk in enumerate(chunks):
            assert isinstance(chunk, DocumentChunk)
            assert chunk.page_number == 1
            assert chunk.chunk_type == "test"
            assert chunk.metadata["chunk_index"] == i
    
    def test_process_pdf(self, processor):
        """Test PDF processing"""
        # Create a simple PDF
        with tempfile.NamedTemporaryFile(delete=False, suffix=".pdf") as tmp:
            from reportlab.pdfgen import canvas
            from reportlab.lib.pagesizes import letter
            
            c = canvas.Canvas(tmp.name, pagesize=letter)
            c.drawString(100, 750, "Test PDF Content")
            c.drawString(100, 700, "This is a test document for processing.")
            c.showPage()
            c.save()
            
            tmp_path = tmp.name
        
        try:
            chunks = processor.process_document(tmp_path, "pdf")
            assert len(chunks) > 0
            assert all(isinstance(chunk, DocumentChunk) for chunk in chunks)
        finally:
            os.unlink(tmp_path)
    
    def test_process_word(self, processor):
        """Test Word document processing"""
        with tempfile.NamedTemporaryFile(delete=False, suffix=".docx") as tmp:
            doc = Document()
            doc.add_heading("Test Document", 0)
            doc.add_paragraph("This is a test paragraph.")
            doc.add_table(rows=2, cols=2)
            doc.save(tmp.name)
            tmp_path = tmp.name
        
        try:
            chunks = processor.process_document(tmp_path, "docx")
            assert len(chunks) > 0
            
            # Check for title chunk
            title_chunks = [c for c in chunks if c.chunk_type == "title"]
            assert len(title_chunks) > 0
            
            # Check for table chunk
            table_chunks = [c for c in chunks if c.chunk_type == "table"]
            assert len(table_chunks) > 0
        finally:
            os.unlink(tmp_path)
    
    def test_process_powerpoint(self, processor):
        """Test PowerPoint processing"""
        with tempfile.NamedTemporaryFile(delete=False, suffix=".pptx") as tmp:
            prs = Presentation()
            slide_layout = prs.slide_layouts[0]
            slide = prs.slides.add_slide(slide_layout)
            title = slide.shapes.title
            subtitle = slide.placeholders[1]
            
            title.text = "Test Presentation"
            subtitle.text = "This is a test slide"
            
            prs.save(tmp.name)
            tmp_path = tmp.name
        
        try:
            chunks = processor.process_document(tmp_path, "pptx")
            assert len(chunks) > 0
            
            # Check for title chunk
            title_chunks = [c for c in chunks if c.chunk_type == "title"]
            assert len(title_chunks) > 0
            assert "Test Presentation" in title_chunks[0].content
        finally:
            os.unlink(tmp_path)
    
    def test_table_formatting(self, processor):
        """Test table formatting"""
        table_data = [
            ["Header 1", "Header 2", "Header 3"],
            ["Row 1 Col 1", "Row 1 Col 2", "Row 1 Col 3"],
            ["Row 2 Col 1", "Row 2 Col 2", "Row 2 Col 3"]
        ]
        
        formatted = processor._format_table(table_data)
        
        # Check markdown table format
        assert "| Header 1 | Header 2 | Header 3 |" in formatted
        assert "|---|---|---|" in formatted
        assert "| Row 1 Col 1 | Row 1 Col 2 | Row 1 Col 3 |" in formatted
    
    def test_unsupported_file_type(self, processor):
        """Test handling of unsupported file types"""
        with pytest.raises(ValueError, match="Unsupported file type"):
            processor.process_document("test.txt", "txt")

if __name__ == "__main__":
    pytest.main([__file__])
