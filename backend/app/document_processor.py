"""Document processing module for PowerPoint and PDF files."""
import os
from typing import Dict, List, Any
from pptx import Presentation
import PyPDF2
from pdf2image import convert_from_path
import pytesseract
from PIL import Image
import asyncio

class DocumentProcessor:
    """Handles document parsing and text extraction."""
    
    def __init__(self):
        self.supported_formats = ['.pptx', '.ppt', '.pdf']
    
    async def process_document(self, file_path: str, document_id: str) -> Dict[str, Any]:
        """Process a document and extract its content."""
        file_extension = os.path.splitext(file_path)[1].lower()
        
        if file_extension in ['.pptx', '.ppt']:
            return await self._process_powerpoint(file_path, document_id)
        elif file_extension == '.pdf':
            return await self._process_pdf(file_path, document_id)
        else:
            raise ValueError(f"Unsupported file format: {file_extension}")
    
    async def _process_powerpoint(self, file_path: str, document_id: str) -> Dict[str, Any]:
        """Extract content from PowerPoint presentations."""
        prs = Presentation(file_path)
        
        extracted_data = {
            "document_id": document_id,
            "type": "powerpoint",
            "slides": [],
            "text": "",
            "page_count": len(prs.slides),
            "tables": [],
            "images": []
        }
        
        all_text = []
        
        for slide_num, slide in enumerate(prs.slides):
            slide_data = {
                "slide_number": slide_num + 1,
                "title": "",
                "content": [],
                "notes": ""
            }
            
            # Extract title
            if slide.shapes.title:
                slide_data["title"] = slide.shapes.title.text
                all_text.append(slide.shapes.title.text)
            
            # Extract content from shapes
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text = shape.text.strip()
                    if text:
                        slide_data["content"].append(text)
                        all_text.append(text)
                
                # Extract tables
                if shape.has_table:
                    table_data = self._extract_table_from_shape(shape)
                    extracted_data["tables"].append({
                        "slide_number": slide_num + 1,
                        "data": table_data
                    })
            
            # Extract notes
            if slide.has_notes_slide:
                notes_text = slide.notes_slide.notes_text_frame.text
                if notes_text:
                    slide_data["notes"] = notes_text
                    all_text.append(f"Notes: {notes_text}")
            
            extracted_data["slides"].append(slide_data)
        
        extracted_data["text"] = "\n\n".join(all_text)
        return extracted_data
    
    async def _process_pdf(self, file_path: str, document_id: str) -> Dict[str, Any]:
        """Extract content from PDF files."""
        extracted_data = {
            "document_id": document_id,
            "type": "pdf",
            "pages": [],
            "text": "",
            "page_count": 0,
            "tables": [],
            "images": []
        }
        
        all_text = []
        
        # Extract text using PyPDF2
        with open(file_path, 'rb') as file:
            pdf_reader = PyPDF2.PdfReader(file)
            extracted_data["page_count"] = len(pdf_reader.pages)
            
            for page_num, page in enumerate(pdf_reader.pages):
                text = page.extract_text()
                page_data = {
                    "page_number": page_num + 1,
                    "text": text
                }
                extracted_data["pages"].append(page_data)
                all_text.append(f"Page {page_num + 1}:\n{text}")
        
        extracted_data["text"] = "\n\n".join(all_text)
        
        # If text extraction is poor, try OCR
        if len(extracted_data["text"].strip()) < 100:
            extracted_data = await self._process_pdf_with_ocr(file_path, document_id, extracted_data)
        
        return extracted_data
    
    async def _process_pdf_with_ocr(self, file_path: str, document_id: str, extracted_data: Dict[str, Any]) -> Dict[str, Any]:
        """Use OCR for PDFs with poor text extraction."""
        try:
            images = convert_from_path(file_path)
            ocr_text = []
            
            for i, image in enumerate(images):
                text = pytesseract.image_to_string(image)
                ocr_text.append(f"Page {i + 1}:\n{text}")
                
                if i < len(extracted_data["pages"]):
                    extracted_data["pages"][i]["text"] = text
            
            extracted_data["text"] = "\n\n".join(ocr_text)
            extracted_data["ocr_used"] = True
            
        except Exception as e:
            print(f"OCR failed: {str(e)}")
            extracted_data["ocr_error"] = str(e)
        
        return extracted_data
    
    def _extract_table_from_shape(self, shape) -> List[List[str]]:
        """Extract table data from PowerPoint shape."""
        table = shape.table
        table_data = []
        
        for row in table.rows:
            row_data = []
            for cell in row.cells:
                row_data.append(cell.text.strip())
            table_data.append(row_data)
        
        return table_data
    
    def _chunk_text(self, text: str, chunk_size: int = 1000, overlap: int = 200) -> List[Dict[str, Any]]:
        """Split text into overlapping chunks for better retrieval."""
        chunks = []
        words = text.split()
        
        for i in range(0, len(words), chunk_size - overlap):
            chunk_words = words[i:i + chunk_size]
            chunk_text = " ".join(chunk_words)
            
            chunks.append({
                "text": chunk_text,
                "start_index": i,
                "end_index": min(i + chunk_size, len(words))
            })
        
        return chunks
