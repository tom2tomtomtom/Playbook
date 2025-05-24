import os
import logging
from typing import List, Dict, Optional
import PyPDF2
import pdfplumber
import fitz  # PyMuPDF
from pptx import Presentation
from docx import Document as DocxDocument
from dataclasses import dataclass
import re
from PIL import Image
import io
import base64

logger = logging.getLogger(__name__)

@dataclass
class DocumentChunk:
    content: str
    page_number: int
    chunk_type: str  # 'text', 'table', 'title', 'image_description'
    metadata: Dict

class DocumentProcessor:
    def __init__(self, chunk_size: int = 500, chunk_overlap: int = 50):
        self.chunk_size = chunk_size
        self.chunk_overlap = chunk_overlap
    
    def process_document(self, file_path: str, file_type: str) -> List[DocumentChunk]:
        """Process a document and return chunks of content"""
        logger.info(f"Processing document: {file_path} (type: {file_type})")
        
        try:
            if file_type in ['pdf']:
                return self.process_pdf_advanced(file_path)
            elif file_type in ['pptx', 'ppt']:
                return self.process_powerpoint(file_path)
            elif file_type in ['docx', 'doc']:
                return self.process_word(file_path)
            else:
                raise ValueError(f"Unsupported file type: {file_type}")
        except Exception as e:
            logger.error(f"Error processing document: {e}")
            raise
    
    def process_pdf_advanced(self, file_path: str) -> List[DocumentChunk]:
        """Advanced PDF processing with table and image extraction"""
        chunks = []
        
        # Try multiple PDF processing methods for robustness
        try:
            # Method 1: Use pdfplumber for better table extraction
            with pdfplumber.open(file_path) as pdf:
                for page_num, page in enumerate(pdf.pages):
                    # Extract text
                    text = page.extract_text()
                    if text and text.strip():
                        text_chunks = self._create_chunks(
                            text, 
                            page_num + 1, 
                            'pdf_text'
                        )
                        chunks.extend(text_chunks)
                    
                    # Extract tables
                    tables = page.extract_tables()
                    for i, table in enumerate(tables):
                        if table:
                            table_text = self._format_table(table)
                            chunks.append(DocumentChunk(
                                content=table_text,
                                page_number=page_num + 1,
                                chunk_type='table',
                                metadata={
                                    'table_index': i,
                                    'rows': len(table),
                                    'columns': len(table[0]) if table else 0
                                }
                            ))
        except Exception as e:
            logger.warning(f"pdfplumber failed, trying PyMuPDF: {e}")
            
            # Method 2: Use PyMuPDF for better image handling
            doc = fitz.open(file_path)
            for page_num in range(len(doc)):
                page = doc[page_num]
                
                # Extract text
                text = page.get_text()
                if text.strip():
                    text_chunks = self._create_chunks(
                        text,
                        page_num + 1,
                        'pdf_text'
                    )
                    chunks.extend(text_chunks)
                
                # Extract images
                image_list = page.get_images()
                for img_index, img in enumerate(image_list):
                    try:
                        # Get image data
                        xref = img[0]
                        pix = fitz.Pixmap(doc, xref)
                        if pix.n - pix.alpha < 4:  # GRAY or RGB
                            img_data = pix.tobytes("png")
                            chunks.append(DocumentChunk(
                                content=f"[Image on page {page_num + 1}]",
                                page_number=page_num + 1,
                                chunk_type='image_description',
                                metadata={
                                    'image_index': img_index,
                                    'width': pix.width,
                                    'height': pix.height,
                                    'has_image': True
                                }
                            ))
                    except Exception as e:
                        logger.warning(f"Error extracting image: {e}")
            
            doc.close()
        
        return chunks
    
    def process_word(self, file_path: str) -> List[DocumentChunk]:
        """Process Word documents"""
        chunks = []
        doc = DocxDocument(file_path)
        
        # Extract paragraphs
        full_text = []
        for para_num, paragraph in enumerate(doc.paragraphs):
            if paragraph.text.strip():
                # Check if it's a heading
                if paragraph.style.name.startswith('Heading'):
                    chunks.append(DocumentChunk(
                        content=paragraph.text,
                        page_number=1,  # Word doesn't have clear page breaks
                        chunk_type='title',
                        metadata={
                            'style': paragraph.style.name,
                            'paragraph_index': para_num
                        }
                    ))
                full_text.append(paragraph.text)
        
        # Create chunks from full text
        if full_text:
            text_chunks = self._create_chunks(
                "\n".join(full_text),
                1,
                'word_text'
            )
            chunks.extend(text_chunks)
        
        # Extract tables
        for table_num, table in enumerate(doc.tables):
            table_data = []
            for row in table.rows:
                row_data = [cell.text.strip() for cell in row.cells]
                table_data.append(row_data)
            
            table_text = self._format_table(table_data)
            chunks.append(DocumentChunk(
                content=table_text,
                page_number=1,
                chunk_type='table',
                metadata={
                    'table_index': table_num,
                    'rows': len(table.rows),
                    'columns': len(table.columns)
                }
            ))
        
        return chunks

    def process_powerpoint(self, file_path: str) -> List[DocumentChunk]:
        """Extract content from PowerPoint files with improved handling"""
        chunks = []
        presentation = Presentation(file_path)
        
        for slide_num, slide in enumerate(presentation.slides):
            slide_content = []
            
            # Extract slide title
            if slide.shapes.title:
                title = slide.shapes.title.text
                if title:
                    chunks.append(DocumentChunk(
                        content=title,
                        page_number=slide_num + 1,
                        chunk_type='title',
                        metadata={'slide_number': slide_num + 1}
                    ))
            
            # Extract text from all shapes
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text:
                    slide_content.append(shape.text)
                
                # Extract table content
                if shape.has_table:
                    table_data = []
                    for row in shape.table.rows:
                        row_data = [cell.text for cell in row.cells]
                        table_data.append(row_data)
                    
                    table_text = self._format_table(table_data)
                    chunks.append(DocumentChunk(
                        content=table_text,
                        page_number=slide_num + 1,
                        chunk_type='table',
                        metadata={
                            'slide_number': slide_num + 1,
                            'rows': len(shape.table.rows),
                            'columns': len(shape.table.columns)
                        }
                    ))
            
            # Extract notes
            if slide.has_notes_slide and slide.notes_slide.notes_text_frame.text:
                slide_content.append(f"Speaker Notes: {slide.notes_slide.notes_text_frame.text}")
            
            # Create chunks from slide content
            if slide_content:
                full_text = "\n".join(slide_content)
                slide_chunks = self._create_chunks(
                    full_text,
                    slide_num + 1,
                    'slide_content'
                )
                chunks.extend(slide_chunks)
        
        return chunks
    
    def _create_chunks(self, text: str, page_number: int, chunk_type: str) -> List[DocumentChunk]:
        """Split text into overlapping chunks with improved handling"""
        chunks = []
        
        # Clean text
        text = re.sub(r'\s+', ' ', text.strip())
        words = text.split()
        
        if len(words) <= self.chunk_size:
            # If text is small enough, return as single chunk
            chunks.append(DocumentChunk(
                content=text,
                page_number=page_number,
                chunk_type=chunk_type,
                metadata={'word_count': len(words)}
            ))
            return chunks
        
        # Create overlapping chunks
        for i in range(0, len(words), self.chunk_size - self.chunk_overlap):
            chunk_words = words[i:i + self.chunk_size]
            chunk_text = " ".join(chunk_words)
            
            chunks.append(DocumentChunk(
                content=chunk_text,
                page_number=page_number,
                chunk_type=chunk_type,
                metadata={
                    'word_count': len(chunk_words),
                    'chunk_index': i // (self.chunk_size - self.chunk_overlap),
                    'start_word': i,
                    'end_word': min(i + self.chunk_size, len(words))
                }
            ))
            
            if i + self.chunk_size >= len(words):
                break
        
        return chunks
    
    def _format_table(self, table_data: List[List[str]]) -> str:
        """Format table data as markdown"""
        if not table_data:
            return ""
        
        # Create markdown table
        lines = []
        
        # Header
        if len(table_data) > 0:
            lines.append("| " + " | ".join(str(cell) for cell in table_data[0]) + " |")
            lines.append("|" + "---|" * len(table_data[0]))
        
        # Rows
        for row in table_data[1:]:
            lines.append("| " + " | ".join(str(cell) for cell in row) + " |")
        
        return "\n".join(lines)
